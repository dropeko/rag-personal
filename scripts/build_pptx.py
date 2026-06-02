"""
build_pptx.py — Refaz o deck Digital Labs com texto enxuto e dados objetivos.

Limites por shape derivados do template (em chars):

  Cover/closing (87pt huge):        ≤ 18 chars/linha · 2 linhas max
  Promessa headline (87pt huge):    ≤ 22 chars total (1-2 linhas)
  Promessa subtitle (21pt):         ≤ 100 chars (~2 linhas em 7.6")
  Visão produto headline (54pt):    ≤ 28 chars (1-2 linhas em 7.77")
  Visão produto subtitle (21pt):    ≤ 110 chars
  3 módulos headline (54pt):        ≤ 24 chars (em 6.10")
  3 cards problema headline (54pt): ≤ 30 chars (em 7.78")
  3 cards problema sub (21pt):      ≤ 80 chars (em 7.78")
  Card 01 (shape 11 estreito):      ≤ 75 chars body (16.5pt)
  Cards 02/03 (shape 18/25):        ≤ 80 chars body (16.5pt)
  Módulo card body (16.5pt):        ≤ 75 chars (em 5.0")
  Casos de uso body (16.5pt):       ≤ 65 chars (em 4.68")
  Antes/depois headline (34.5pt):   ≤ 16 chars total
  Antes/depois col headline (28.5pt): ≤ 35 chars
  Antes/depois col body (16.5pt):   ≤ 130 chars
  4 fases headline (34.5pt):        ≤ 26 chars
  4 fases title (23pt):             ≤ 11 chars
  4 fases body (16.5pt):            ≤ 55 chars
  Pull quote (52.5pt):              ≤ 100 chars
  Prova huge (100.5pt):             ≤ 5 chars
  Prova subtitle (22.5pt):          ≤ 80 chars total
"""
from copy import deepcopy
from pathlib import Path
from pptx import Presentation

TEMPLATE = "/tmp/template.pptx"
OUT = Path(__file__).resolve().parent.parent / "references" / "Apresentacao_RAG_2026-06-01.pptx"

prs = Presentation(TEMPLATE)
TEMPLATE_SLIDES = list(prs.slides)
TOTAL = 17


def clone_slide(src_slide):
    new_slide = prs.slides.add_slide(src_slide.slide_layout)
    for shp in list(new_slide.shapes):
        shp._element.getparent().remove(shp._element)
    for shp in src_slide.shapes:
        new_slide.shapes._spTree.append(deepcopy(shp._element))
    return new_slide


def find_shape(slide, name):
    for sh in slide.shapes:
        if sh.name == name:
            return sh
    return None


def set_paragraphs(shape, lines):
    """Preserva formatação; deixa parágrafos vazios verdadeiramente vazios."""
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    paras = list(tf.paragraphs)
    for i, p in enumerate(paras):
        if i < len(lines):
            new = lines[i]
            if p.runs:
                p.runs[0].text = new
                for r in p.runs[1:]:
                    r.text = ""
            else:
                if new:
                    p.text = new
        else:
            for r in p.runs:
                r.text = ""


def set_table_cells(shape, mapping):
    if shape.shape_type != 19:
        return
    tbl = shape.table
    for (r, c), text in mapping.items():
        cell = tbl.cell(r, c)
        tf = cell.text_frame
        paras = list(tf.paragraphs)
        if paras and paras[0].runs:
            paras[0].runs[0].text = text
            for run in paras[0].runs[1:]:
                run.text = ""
            for p in paras[1:]:
                for run in p.runs:
                    run.text = ""
        else:
            cell.text = text


def reorder_slides(prs, slides_in_order):
    ns = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"
    sldIdLst = prs.slides._sldIdLst
    by_partname = {}
    for sldId in list(sldIdLst):
        rId = sldId.attrib[f"{ns}id"]
        part = prs.part.related_part(rId)
        by_partname[part.partname] = sldId
        sldIdLst.remove(sldId)
    for sl in slides_in_order:
        sldIdLst.append(by_partname[sl.part.partname])


# ============================================================================
# 17 builders — texto enxuto, dados objetivos
# ============================================================================

def b_cover(s):
    set_paragraphs(find_shape(s, "object 3"), ["ISQ BRASIL  /  NIT-DEV"])
    set_paragraphs(find_shape(s, "object 4"), ["JUN 2026"])
    set_paragraphs(find_shape(s, "object 5"), ["ANÁLISE TÉCNICA — SPRINT"])
    set_paragraphs(find_shape(s, "object 6"), [
        "Stack de RAG",
        "para inspeção.",
        "ChromaDB local vs. OpenAI Embeddings. "
        "Embedder, vector store, LLM e ingestão contínua.",
    ])
    set_paragraphs(find_shape(s, "object 9"), ["ANÁLISE  /  16 : 9"])


def b_roteiro(s):
    set_paragraphs(find_shape(s, "object 2"), ["ROTEIRO"])
    set_paragraphs(find_shape(s, "object 3"), [
        "Do corpus à recomendação.",
    ])
    set_table_cells(find_shape(s, "object 5"), {
        (0, 0): "01\tContexto e descobertas",
        (0, 2): "\n02",
        (0, 3): "Stack atual e bug detectado",
        (1, 0): "03\tCorpus e ingestão contínua",
        (1, 2): "\n04",
        (1, 3): "Custos e latência",
        (2, 0): "05\tBench parcial — resultados",
        (2, 2): "\n06",
        (2, 3): "Proposta e próximos passos",
    })
    set_paragraphs(find_shape(s, "object 9"), [f"02 / {TOTAL}"])


def b_sumario(s):
    # promessa central — 87pt headline curto + 21pt sub curto
    set_paragraphs(find_shape(s, "object 3"), [
        "SUMÁRIO EXECUTIVO",
        "",
        "",
        "Precisão, não escala.",  # 21 chars
        "Corpus de 8.021 chunks. Bug atual entre indexação e consulta. "
        "Embedding custa centavos; a decisão está no retrieval e no LLM gerador.",  # 132 chars (~3 linhas)
    ])
    set_paragraphs(find_shape(s, "object 6"), [f"03 / {TOTAL}"])


def b_descobertas(s):
    set_paragraphs(find_shape(s, "object 2"), ["ACHADOS"])
    set_paragraphs(find_shape(s, "object 3"), [
        "Três achados que",            # 16 chars 54pt
        "mudam a sprint.",             # 15 chars 54pt
        "Direto do chroma.sqlite3 e do código do db-backend.",  # 51 chars 21pt
        "",
    ])
    set_paragraphs(find_shape(s, "object 10"), ["01", "Mismatch"])
    set_paragraphs(find_shape(s, "object 11"), [
        "Indexação usa fine-tunado; consulta usa MPNet vanilla.",  # 54 chars
    ])
    set_paragraphs(find_shape(s, "object 18"), [
        "02", "Corpus pequeno",
        "78 laudos e 7.943 captions de metalografia. Chunk médio: 400 tokens.",  # 67 chars
    ])
    set_paragraphs(find_shape(s, "object 25"), [
        "03", "Ingestão contínua",
        "Novos laudos precisam de pipeline assíncrono — ainda não desenhado.",  # 67 chars
    ])
    set_paragraphs(find_shape(s, "object 28"), [f"04 / {TOTAL}"])


def b_stack_atual(s):
    # visão do produto — 54pt headline + 21pt sub
    set_paragraphs(find_shape(s, "object 2"), ["STACK ATUAL"])
    set_paragraphs(find_shape(s, "object 3"), [
        "Chroma local + fine-tunado interno.",   # 35 chars 54pt → 2-3 linhas
        "NestJS via chromadb (JS). Index com fine-tunado (ml_wo); "
        "consultas com MPNet vanilla.",          # 85 chars 21pt
    ])
    set_paragraphs(find_shape(s, "object 8"), ["STACK / ml_wo"])
    set_paragraphs(find_shape(s, "object 9"), ["V 1"])
    set_paragraphs(find_shape(s, "object 11"), [
        "CAMADAS  ·  EMBEDDER  ·  STORE  ·  CLIENTE  ·  LLM",
    ])
    set_paragraphs(find_shape(s, "object 14"), [f"05 / {TOTAL}"])


def b_bug_detalhe(s):
    set_paragraphs(find_shape(s, "object 2"), ["MISMATCH"])
    set_paragraphs(find_shape(s, "object 3"), [
        "Índice e consulta",              # 17 chars
        "usam modelos diferentes.",       # 24 chars
        "Resultado: espaços vetoriais incompatíveis em produção.",  # 56 chars
        "",
    ])
    set_paragraphs(find_shape(s, "object 10"), ["01", "Indexação"])
    set_paragraphs(find_shape(s, "object 11"), [
        "fine_tuned_report_model — 768 dim. Fora do repositório.",  # 56 chars
    ])
    set_paragraphs(find_shape(s, "object 18"), [
        "02", "Consulta",
        "Xenova/MPNet vanilla em ONNX — 768 dim. Aplicado em cada chamada.",  # 65 chars
    ])
    set_paragraphs(find_shape(s, "object 25"), [
        "03", "Consequência",
        "Recall degradado. Magnitude só será medida no bench neural.",  # 60 chars
    ])
    set_paragraphs(find_shape(s, "object 28"), [f"06 / {TOTAL}"])


def b_corpus_real(s):
    # antes/depois — headline curto + 2 colunas
    set_paragraphs(find_shape(s, "object 2"), ["CORPUS"])
    set_paragraphs(find_shape(s, "object 3"), [
        "78 laudos",          # 9 chars 34.5pt
        "+ 7.943 captions.",  # 17 chars 34.5pt
    ])
    set_paragraphs(find_shape(s, "object 7"), ["report_texts"])
    set_paragraphs(find_shape(s, "object 8"), [
        "78 chunks — um por laudo.",       # 25 chars 28.5pt
        "Seção “Discussão dos Resultados”. ~400 tokens por chunk. "
        "100 % Análise de Falha PETROBRÁS.",   # 90 chars 16.5pt
    ])
    set_paragraphs(find_shape(s, "object 12"), ["captions"])
    set_paragraphs(find_shape(s, "object 13"), [
        "7.943 captions",          # 14 chars 28.5pt
        "de metalografia.",        # 16 chars 28.5pt
        "~100 captions por laudo. Atende consultas onde a resposta "
        "está na imagem, não no texto.",   # 87 chars 16.5pt
    ])
    set_paragraphs(find_shape(s, "object 16"), [f"07 / {TOTAL}"])


def b_ingestao(s):
    # 4 fases — headline curto + 4 cards
    set_paragraphs(find_shape(s, "object 2"), ["INGESTÃO CONTÍNUA"])
    set_paragraphs(find_shape(s, "object 3"), [
        "Novo laudo buscável em < 1 minuto.",   # 34 chars 34.5pt
    ])
    set_paragraphs(find_shape(s, "object 5"), ["01"])
    set_paragraphs(find_shape(s, "object 6"), [
        "Gerar",
        "report-generate finaliza o DOCX e emite evento.",   # 47 chars
    ])
    set_paragraphs(find_shape(s, "object 8"), ["02"])
    set_paragraphs(find_shape(s, "object 9"), [
        "Enfileirar",
        "Fila no Postgres existente (pg-boss).",   # 37 chars
    ])
    set_paragraphs(find_shape(s, "object 11"), ["03"])
    set_paragraphs(find_shape(s, "object 12"), [
        "Indexar",
        "Worker embeda e faz upsert idempotente.",   # 39 chars
    ])
    set_paragraphs(find_shape(s, "object 14"), ["04"])
    set_paragraphs(find_shape(s, "object 15"), [
        "Servir",
        "Coleção atualizada; busca enxerga o laudo.",   # 42 chars
    ])
    set_paragraphs(find_shape(s, "object 18"), [f"08 / {TOTAL}"])


def b_opcoes(s):
    # arquitetura de valor — headline curto + 3 módulos
    set_paragraphs(find_shape(s, "object 2"), ["OPÇÕES"])
    set_paragraphs(find_shape(s, "object 3"), [
        "Três caminhos,",     # 14 chars 54pt
        "mesmo bench.",       # 12 chars 54pt
    ])
    set_paragraphs(find_shape(s, "object 10"), [
        "Opção A", "Status quo coerente",
        "Reindex com vanilla. Chroma local. Corrige o bug sem migrar.",  # 60 chars
    ])
    set_paragraphs(find_shape(s, "object 17"), [
        "Opção B", "OpenAI Embeddings",
        "text-embedding-3-small @768 via API. A/B sem refazer schema.",  # 61 chars
    ])
    set_paragraphs(find_shape(s, "object 24"), [
        "Opção C", "Híbrida",
        "OpenAI + pgvector no Postgres. Sem chroma-server Python.",  # 56 chars
    ])
    set_paragraphs(find_shape(s, "object 27"), [f"09 / {TOTAL}"])


def b_custos(s):
    # casos de uso — headline curto + 3 cards com números
    set_paragraphs(find_shape(s, "object 2"), ["CUSTOS"])
    set_paragraphs(find_shape(s, "object 3"), [
        "Embedding ≠",       # 11 chars 54pt
        "custo real.",       # 11 chars 54pt
    ])
    set_paragraphs(find_shape(s, "object 10"), [
        "Item 01", "Indexação",
        "Corpus completo: US$ 0,06 Standard · US$ 0,03 Batch.",  # 53 chars
    ])
    set_paragraphs(find_shape(s, "object 17"), [
        "Item 02", "Consulta",
        "US$ 0,0000006 por chamada. 100 k/mês = US$ 0,06.",  # 48 chars
    ])
    set_paragraphs(find_shape(s, "object 24"), [
        "Item 03", "LLM gpt-4o-mini",
        "US$ 0,0008 por resposta. 10 k/mês = US$ 8.",  # 42 chars
    ])
    set_paragraphs(find_shape(s, "object 27"), [f"10 / {TOTAL}"])


def b_latencia(s):
    # antes/depois — números objetivos lado a lado
    set_paragraphs(find_shape(s, "object 2"), ["LATÊNCIA"])
    set_paragraphs(find_shape(s, "object 3"), [
        "Embedder estimado;",   # 18 chars 34.5pt
        "store medido.",        # 13 chars 34.5pt
    ])
    set_paragraphs(find_shape(s, "object 7"), ["Estimativa"])
    set_paragraphs(find_shape(s, "object 8"), [
        "Embedder e LLM dominam.",      # 23 chars 28.5pt
        "1 query: 80–250 ms local · 60–180 ms OpenAI. "
        "LLM 4o-mini: 1,5–4 s por resposta.",   # 80 chars 16.5pt
    ])
    set_paragraphs(find_shape(s, "object 12"), ["Medido 2026-06-01"])
    set_paragraphs(find_shape(s, "object 13"), [
        "Chroma HNSW",                  # 11 chars 28.5pt
        "não é gargalo.",               # 14 chars 28.5pt
        "p50: 4,58 ms (78 docs) · 5,08 ms (7.943 docs). "
        "p99: 5,99 ms · 6,97 ms.",   # 70 chars 16.5pt
    ])
    set_paragraphs(find_shape(s, "object 16"), [f"11 / {TOTAL}"])


def b_tfidf(s):
    # prova e evidência — número único enorme
    set_paragraphs(find_shape(s, "object 2"), [
        "BENCH — 2026-06-01",   # 18 chars 11pt
        "67 %",                  # 4 chars 100.5pt
        "Recall@5 do baseline TF-IDF — piso a superar.",  # 47 chars 22.5pt → 2 linhas
        "234 consultas, 78 chunks.",                       # 25 chars 22.5pt → 1 linha
    ])
    set_paragraphs(find_shape(s, "object 5"), [f"12 / {TOTAL}"])


def b_stack_proposta(s):
    # arquitetura — 3 decisões
    set_paragraphs(find_shape(s, "object 2"), ["PROPOSTA"])
    set_paragraphs(find_shape(s, "object 3"), [
        "Três decisões",      # 13 chars 54pt
        "para a stack.",      # 13 chars 54pt
    ])
    set_paragraphs(find_shape(s, "object 10"), [
        "Decisão 1", "Embedder",
        "text-embedding-3-small @768. Xenova como fallback.",  # 50 chars
    ])
    set_paragraphs(find_shape(s, "object 17"), [
        "Decisão 2", "Vector store",
        "Chroma agora; pgvector no Postgres entra no roadmap.",  # 53 chars
    ])
    set_paragraphs(find_shape(s, "object 24"), [
        "Decisão 3", "LLM gerador",
        "gpt-4o-mini baseline. Subir só com sinal de regressão.",  # 55 chars
    ])
    set_paragraphs(find_shape(s, "object 27"), [f"13 / {TOTAL}"])


def b_criterios(s):
    # pull quote 52.5pt — máximo ~100 chars
    set_paragraphs(find_shape(s, "object 2"), ["CRITÉRIOS"])
    set_paragraphs(find_shape(s, "object 3"), [
        "“Migrar para OpenAI somente se Recall@5 ≥ +5 pp, "
        "MRR ≥ +0,05 e p95 ≤ 400 ms.”",   # 80 chars 52.5pt
        "DPA · ZDR · MASCARAMENTO DE PII COMO PRÉ-CONDIÇÕES",  # 50 chars 13.5pt
    ])
    set_paragraphs(find_shape(s, "object 6"), [f"14 / {TOTAL}"])


def b_restricao(s):
    set_paragraphs(find_shape(s, "object 2"), ["BLOQUEIO"])
    set_paragraphs(find_shape(s, "object 3"), [
        "A TI ainda não",        # 15 chars 54pt
        "liberou a rede.",       # 15 chars 54pt
        "Bench parcial entregue; bench neural depende de liberação.",  # 60 chars 21pt
        "",
    ])
    set_paragraphs(find_shape(s, "object 10"), ["01", "huggingface.co"])
    set_paragraphs(find_shape(s, "object 11"), [
        "MPNet, BGE-m3 e demais open source. Liberação pontual.",  # 56 chars
    ])
    set_paragraphs(find_shape(s, "object 18"), [
        "02", "api.openai.com",
        "Embeddings v3 e LLM gerador. Liberação permanente, allowlist.",  # 61 chars
    ])
    set_paragraphs(find_shape(s, "object 25"), [
        "03", "Plano B",
        "Bench em estação pessoal. Scripts reproduzíveis em ~3 horas.",  # 60 chars
    ])
    set_paragraphs(find_shape(s, "object 28"), [f"15 / {TOTAL}"])


def b_proximos_passos(s):
    # 4 fases — headline curto + 4 fases compactas
    set_paragraphs(find_shape(s, "object 2"), ["PRÓXIMOS PASSOS"])
    set_paragraphs(find_shape(s, "object 3"), [
        "Quatro frentes em paralelo.",   # 27 chars 34.5pt
    ])
    set_paragraphs(find_shape(s, "object 5"), ["01"])
    set_paragraphs(find_shape(s, "object 6"), [
        "Liberar",
        "Chamado TI: huggingface.co + api.openai.com.",   # 44 chars
    ])
    set_paragraphs(find_shape(s, "object 8"), ["02"])
    set_paragraphs(find_shape(s, "object 9"), [
        "Rodar",
        "prepare.py + run-queries + score.py com chave OpenAI.",   # 53 chars
    ])
    set_paragraphs(find_shape(s, "object 11"), ["03"])
    set_paragraphs(find_shape(s, "object 12"), [
        "Decidir",
        "Destino do fine_tuned_report_model.",   # 35 chars
    ])
    set_paragraphs(find_shape(s, "object 14"), ["04"])
    set_paragraphs(find_shape(s, "object 15"), [
        "Implantar",
        "Pipeline de ingestão sobre o Postgres existente.",   # 48 chars
    ])
    set_paragraphs(find_shape(s, "object 18"), [f"16 / {TOTAL}"])


def b_proxima_decisao(s):
    # closing — 87pt headline curto
    set_paragraphs(find_shape(s, "object 3"), ["NIT-DEV  /  DROPEKO"])
    set_paragraphs(find_shape(s, "object 4"), ["ISQ DIGITAL LABS"])
    set_paragraphs(find_shape(s, "object 5"), [
        "PRÓXIMA DECISÃO",
        "",
        "",
        "Liberar a TI.",      # 13 chars 87pt — 1-2 linhas
        "Com a rede liberada, o bench neural sai em uma tarde. "
        "Sem ela, a stack fica travada no piso do TF-IDF.",   # 100 chars 21pt
    ])
    set_paragraphs(find_shape(s, "object 8"), [f"17 / {TOTAL}"])


# ---------- plano de slides ----------
PLAN = [
    ("inplace", 0, b_cover),
    ("inplace", 1, b_roteiro),
    ("inplace", 2, b_sumario),
    ("inplace", 3, b_descobertas),
    ("inplace", 4, b_stack_atual),
    ("clone",   3, b_bug_detalhe),
    ("inplace", 9, b_corpus_real),
    ("inplace", 10, b_ingestao),
    ("inplace", 5, b_opcoes),
    ("inplace", 7, b_custos),
    ("clone",   9, b_latencia),
    ("inplace", 8, b_tfidf),
    ("clone",   5, b_stack_proposta),
    ("inplace", 6, b_criterios),
    ("clone",   3, b_restricao),
    ("clone",   10, b_proximos_passos),
    ("inplace", 11, b_proxima_decisao),
]

ordered_slides = []
for mode, idx, builder in PLAN:
    sl = TEMPLATE_SLIDES[idx] if mode == "inplace" else clone_slide(TEMPLATE_SLIDES[idx])
    builder(sl)
    ordered_slides.append(sl)

reorder_slides(prs, ordered_slides)

OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(OUT)
print(f"deck gerado: {OUT}  ({len(list(prs.slides))} slides)")
